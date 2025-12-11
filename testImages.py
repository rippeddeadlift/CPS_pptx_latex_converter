from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Replace with your actual path
PPT_PATH = "./input/Algorithmik.pptx" 

prs = Presentation(PPT_PATH)
print(f"--- Scanning {len(prs.slides)} slides ---")

for i, slide in enumerate(prs.slides):
    print(f"\nSlide {i+1}:")
    for shape in slide.shapes:
        # Print the Type and Name of every object
        print(f"  - Shape: '{shape.name}' | Type: {shape.shape_type}")
        
        # Check if it has a fill picture (common for fancy slides)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if hasattr(shape.fill, 'type') and shape.fill.type == 6: # 6 = Picture Fill
                    print(f"    *** HIDDEN IMAGE FOUND (AutoShape with Picture Fill) ***")
            except:
                pass
                
        # Check if it is a Group (images often hide inside groups)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
             print(f"    *** GROUP FOUND (Images might be inside) ***")