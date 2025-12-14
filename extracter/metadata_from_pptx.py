from pptx.enum.shapes import PP_PLACEHOLDER

def get_institute_heuristic(prs, known_title, known_author):
    """
    Holt den Text direkt aus dem offiziellen Footer-Placeholder des Slide Masters.
    Das ist der sauberste Weg für wiederkehrende Texte wie Institutsnamen.
    """
    try:
        if not prs.slides:
            return ""

        # 1. Zugriff auf den Master der ersten Folie
        first_slide = prs.slides[0]
        master = first_slide.slide_layout.slide_master

        # 2. Durchsuche NUR die offiziellen Platzhalter im Master
        for shape in master.placeholders:
            # Wir prüfen exakt auf den Typ FOOTER (Enum ID 15)
            if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                text = shape.text.strip()
                if text:
                    return text

        # 3. Fallback: Manchmal ist der Footer kein "Placeholder", sondern eine Textbox im Master.
        # Wir suchen nach einer Textbox ganz unten im Master (untere 10%).
        slide_height = prs.slide_height
        for shape in master.shapes:
            if not shape.has_text_frame: 
                continue
            
            # Ist es weit genug unten? (Untere 15% der Seite)
            if shape.top > (slide_height * 0.85):
                text = shape.text.strip()
                # Filter: Keine Seitenzahlen ("<#>") oder Datum
                if text and not text.isdigit() and len(text) > 3:
                     # Filter: Schließe typische Platzhalter-Texte aus
                    if "datum" not in text.lower() and "date" not in text.lower():
                         return text

    except Exception as e:
        print(f"   [WARN] Error reading master footer: {e}")

    return ""