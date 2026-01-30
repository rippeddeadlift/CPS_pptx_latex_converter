import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_media_from_pptx(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation(pptx_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    layout_data_by_slide = {}
    global_media_count = 1 


    for i, slide in enumerate(prs.slides):
        slide_index = i
        slide_media = []
        
        rels = slide.part.rels

        for shape in slide.shapes:
            global_media_count = _process_shape(
                shape, rels, slide_media, output_dir, global_media_count, slide_width, slide_height
            )

        if slide_media:
            layout_data_by_slide[slide_index] = slide_media
            print(f"Slide {i+1}: Found {len(slide_media)} media items")

    return layout_data_by_slide

def _process_shape(shape, rels, slide_media, output_dir, count, s_width, s_height):
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            count = _process_shape(child, rels, slide_media, output_dir, count, s_width, s_height)
        return count

    xml_text = shape.element.xml
    rids = re.findall(r'r:embed="([^"]+)"', xml_text) + \
           re.findall(r'r:link="([^"]+)"', xml_text)

    for rid in rids:
        if rid in rels:
            rel = rels[rid]
            try:
                if rel.is_external: continue
                
                part = rel.target_part
                ctype = part.content_type.lower()

                if ctype.startswith('video/') or ctype in ['application/x-mplayer2']:
                    ext = "mp4"
                    if "wmv" in ctype: ext = "wmv"
                    elif "avi" in ctype: ext = "avi"
                    elif "quicktime" in ctype: ext = "mov"
                    
                    video_filename = f"media_{count}.{ext}"
                    video_filepath = os.path.join(output_dir, video_filename)
                    
                    with open(video_filepath, "wb") as f:
                        f.write(part.blob)
                    
                    poster_path_json = ""
                    try:
                        if hasattr(shape, "image"):
                            img = shape.image
                            img_ext = img.ext
                            poster_filename = f"media_{count}_poster.{img_ext}"
                            poster_filepath = os.path.join(output_dir, poster_filename)
                            
                            with open(poster_filepath, "wb") as f_img:
                                f_img.write(img.blob)
                            
                            parent_folder = Path(poster_filepath).parent.name
                            poster_path_json = f"{parent_folder}/{poster_filename}"
                            print(f"      [POSTER] Saved preview image: {poster_filename}")
                    except Exception as e:
                        print(f"      [INFO] No poster image extracted: {e}")

                    print(f"      [HIT] Video found: {video_filename}")
                    
                    relative_folder = Path(video_filepath).parent.name
                    video_json_path = f"{relative_folder}/{video_filename}"
                    
                    left = shape.left / s_width
                    top = shape.top / s_height
                    width = shape.width / s_width
                    height = shape.height / s_height

                    slide_media.append({
                        "type": "video",
                        "filename": video_filename,
                        "path": video_json_path,
                        "poster_path": poster_path_json, 
                        "geometry": {
                            "x": round(left, 3), "y": round(top, 3),
                            "w": round(width, 3), "h": round(height, 3)
                        }
                    })
                    return count + 1

            except Exception as e:
                continue

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return _save_shape_image(shape, slide_media, output_dir, count, s_width, s_height)

    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
         if hasattr(shape, 'image') and shape.image:
            return _save_shape_image(shape, slide_media, output_dir, count, s_width, s_height)
            
    return count

def _save_shape_image(shape, slide_media, output_dir, count, s_width, s_height):
    try:
        image = shape.image
        ext = image.ext
        filename = f"image_{count}.{ext}"
        filepath = os.path.join(output_dir, filename)
        
        with open(filepath, "wb") as f:
            f.write(image.blob)
        
        _append_to_list(slide_media, "picture", filename, filepath, shape, s_width, s_height)
        return count + 1
    except:
        return count

def _append_to_list(slide_media, type_name, filename, full_path, shape, s_width, s_height):
    relative_folder_name = Path(full_path).parent.name 
    json_relative_path = f"{relative_folder_name}/{filename}"
    
    left = shape.left / s_width
    top = shape.top / s_height
    width = shape.width / s_width
    height = shape.height / s_height
    # --- DEBUG PRINT ANFANG ---
    print(f"      🔍 [GEOMETRY CHECK] '{filename}'")
    print(f"          -> X: {left:.3f} ({(left*100):.1f}%) | Y: {top:.3f} ({(top*100):.1f}%)")
    print(f"          -> W: {width:.3f} ({(width*100):.1f}%) | H: {height:.3f} ({(height*100):.1f}%)")
    # --- DEBUG PRINT ENDE ---
    slide_media.append({
        "type": type_name,
        "filename": filename,
        "path": json_relative_path,
        "geometry": {
            "x": round(left, 3),
            "y": round(top, 3),
            "w": round(width, 3),
            "h": round(height, 3)
        }
    })