from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from helper.utils import YELLOW, RESET, BLUE

def get_institute_heuristic(prs) -> str:
    """
    Attempts to retrieve the institute or organization name from the Slide Master footer.

    First checks for an explicit Footer placeholder in the master slide. If not found,
    scans for text shapes located in the bottom 15% of the master slide, filtering out
    numbers (likely page numbers) and date-related strings.
    """
    try:
        if not prs.slides:
            return ""

        first_slide = prs.slides[0]
        master = first_slide.slide_layout.slide_master

        for shape in master.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                text = shape.text.strip()
                if text:
                    return text

        slide_height = prs.slide_height
        bottom_threshold = slide_height * 0.85

        for shape in master.shapes:
            if not shape.has_text_frame: 
                continue
            
            if shape.top > bottom_threshold:
                text = shape.text.strip()
                
                if text and not text.isdigit() and len(text) > 3:
                    text_lower = text.lower()
                    if "datum" not in text_lower and "date" not in text_lower:
                         return text

    except Exception as e:
        print(f"{YELLOW}   [WARN] Error reading master footer: {e}{RESET}")

    return ""

def extract_metadata(config) -> dict:
    """
    Extracts metadata from the PPTX or falls back to defaults.
    """
    try:
        print(f"{BLUE}Extracting PPTX metadata...{RESET}")
        
        prs = Presentation(config.PPTX_INPUT)
        props = prs.core_properties        
        title_text = props.title if props.title else config.PPTX_INPUT.stem        
        author_text = props.author if props.author else "AI Converter"        
        institute_text = props.category if props.category else ""
        
        if not institute_text:
            print(f"{YELLOW}   -> No metadata 'category' found. Trying to guess from Slide Master...{RESET}")
            institute_text = get_institute_heuristic(prs)
            
        if institute_text:
            institute_text = institute_text.replace('\n', r' \\ ')

        meta = {
            "title": title_text,
            "author": author_text,
            "date": props.created.strftime("%d.%m.%Y") if props.created else r"\today",
            "institute": institute_text
        }
        
        return meta

    except Exception as e:
        print(f"{YELLOW}ERROR: Failed to extract metadata: {e}{RESET}")
        return {
            "title": config.PPTX_INPUT.stem,
            "author": "Unknown",
            "date": r"\today",
            "institute": ""
        }