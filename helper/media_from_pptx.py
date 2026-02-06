import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from helper.utils import GREEN, RESET


def extract_media_from_pptx(pptx_path: str, output_dir: str) -> dict:
    """
    Iterates through all slides in the presentation to extract images and videos.

    Creates the output directory if it doesn't exist. For each slide, it processes all shapes
    (including nested groups) to save media files and collect their metadata (geometry, path).
    Returns a dictionary mapping slide indices to lists of extracted media items.
    """
    out_path = Path(output_dir)
    out_path.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    layout_data_by_slide = {}
    global_media_count = 1 

    for i, slide in enumerate(prs.slides):
        slide_media = []
        rels = slide.part.rels

        for shape in slide.shapes:
            global_media_count = _process_shape(
                shape, rels, slide_media, str(out_path), global_media_count, slide_width, slide_height
            )

        if slide_media:
            layout_data_by_slide[i] = slide_media
            print(f"Slide {i+1}: Found {len(slide_media)} media items")

    return layout_data_by_slide

def _process_shape(shape, rels, slide_media: list, output_dir: str, count: int, s_width: float, s_height: float) -> int:
    """
    Recursively processes a shape to extract images or embedded videos.

    Checks if the shape is a group (recurses), a picture (extracts), or contains 
    embedded media relationships (extracts video + poster). Updates the slide_media 
    list with metadata and geometry for any found assets.
    """
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            count = _process_shape(child, rels, slide_media, output_dir, count, s_width, s_height)
        return count

    xml_text = shape.element.xml
    rids = re.findall(r'r:embed="([^"]+)"', xml_text) + re.findall(r'r:link="([^"]+)"', xml_text)

    for rid in rids:
        if rid not in rels: continue
        rel = rels[rid]
        if rel.is_external: continue
        
        try:
            part = rel.target_part
            ctype = part.content_type.lower()

            if ctype.startswith('video/') or ctype == 'application/x-mplayer2':
                ext = "mp4"
                if "wmv" in ctype: ext = "wmv"
                elif "avi" in ctype: ext = "avi"
                elif "quicktime" in ctype: ext = "mov"
                
                video_filename = f"media_{count}.{ext}"
                output_path = Path(output_dir)
                video_filepath = output_path / video_filename
                
                with open(video_filepath, "wb") as f:
                    f.write(part.blob)
                
                poster_path_json = ""
                try:
                    if hasattr(shape, "image") and shape.image:
                        img = shape.image
                        poster_filename = f"media_{count}_poster.{img.ext}"
                        poster_filepath = output_path / poster_filename
                        
                        with open(poster_filepath, "wb") as f_img:
                            f_img.write(img.blob)
                        
                        poster_path_json = f"{poster_filepath.parent.name}/{poster_filename}"
                except Exception:
                    pass 

                left = shape.left / s_width
                top = shape.top / s_height
                width = shape.width / s_width
                height = shape.height / s_height

                video_json_path = f"{video_filepath.parent.name}/{video_filename}"

                slide_media.append({
                    "type": "video",
                    "filename": video_filename,
                    "path": video_json_path,
                    "poster_path": poster_path_json, 
                    "geometry": {
                        "x": round(left, 3), "y": round(top, 3),
                        "w": round(width, 3), "h": round(height, 3)
                    }
                })
                
                print(f"{GREEN}      [HIT] Video extracted: {video_filename}{RESET}")
                return count + 1

        except Exception:
            continue

    is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    is_placeholder_img = (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and 
                          hasattr(shape, 'image') and shape.image)

    if is_picture or is_placeholder_img:
        return _save_shape_image(shape, slide_media, output_dir, count, s_width, s_height)
            
    return count


def _save_shape_image(shape, slide_media: list, output_dir: str, count: int, s_width: float, s_height: float) -> int:
    """
    Extracts and saves an image from a shape to the output directory.

    Writes the binary image data to a file with a sequential filename. If successful, 
    it registers the image in the slide's media list using the helper function. 
    Returns the incremented counter on success, or the original counter if extraction fails.
    """
    try:
        image = shape.image
        ext = image.ext
        filename = f"image_{count}.{ext}"
        filepath = os.path.join(output_dir, filename)
        
        with open(filepath, "wb") as f:
            f.write(image.blob)
        
        _append_to_list(slide_media, "picture", filename, filepath, shape, s_width, s_height)
        return count + 1
    except Exception:
        return count


def _append_to_list(slide_media: list, type_name: str, filename: str, full_path: str, shape, s_width: float, s_height: float) -> None:
    """
    Helper function to normalize media geometry and append metadata to the slide list.

    Calculates the relative position and size (0.0-1.0) of a media shape against the 
    total slide dimensions. It constructs a metadata dictionary containing the 
    relative file path and rounded geometry coordinates, then adds it to the 
    provided media list.
    """
    relative_folder_name = Path(full_path).parent.name 
    json_relative_path = f"{relative_folder_name}/{filename}"
    
    left = shape.left / s_width
    top = shape.top / s_height
    width = shape.width / s_width
    height = shape.height / s_height

    slide_media.append({
        "type": type_name,
        "filename": filename,
        "path": json_relative_path,
        "geometry": {
            "x": round(left, 3),
            "y": round(top, 3),
            "w": round(width, 3),
            "h": round(height, 3)
        }
    })