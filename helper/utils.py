
import subprocess
from pathlib import Path
import sys, re, json
from collections import Counter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

RESET = "\033[0m"
RED = "\033[31m"
GREEN = "\033[32m"
YELLOW = "\033[33m"
BLUE = "\033[34m"

def compile_tex_to_pdf(tex_filename: str | Path, working_dir: str | Path) -> bool:
    """
    Compiles a .tex file to PDF using pdflatex within a specified working directory.

    Prioritizes file existence over the compiler's exit code to determine success.
    Removes any existing PDF before compilation to prevent false positives, then
    checks if a new PDF is generated even if pdflatex reports warnings/errors.
    """
    tex_path = Path(tex_filename)
    work_dir = Path(working_dir)
    file_name = tex_path.name
    pdf_name = tex_path.stem + ".pdf"
    pdf_full_path = work_dir / pdf_name

    try:
        if pdf_full_path.exists():
            pdf_full_path.unlink()
    except OSError:
        pass  

    print(f"{BLUE}Compiling {file_name} in {work_dir}...{RESET}")

    command = [
        "pdflatex",
        "-interaction=nonstopmode",
        file_name
    ]

    try:
        result = subprocess.run(
            command,
            cwd=work_dir,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            check=False
        )

        if pdf_full_path.exists():
            if result.returncode == 0:
                print(f"{GREEN}SUCCESS: PDF generated at {pdf_full_path}{RESET}")
            else:
                print(f"{YELLOW}PDF generated with LaTeX issues at {pdf_full_path}{RESET}")
            return True
        else:
            print(f"{RED}ERROR: PDF compilation failed. No output file.{RESET}")
            print(f"{YELLOW}--- LaTeX Error Log (Last 20 lines) ---{RESET}")
            if result.stdout:
                lines = result.stdout.splitlines()
                print("\n".join(lines[-20:]))
            return False

    except FileNotFoundError:
        print(f"{RED}CRITICAL ERROR: 'pdflatex' not found.{RESET}")
        print(f"{YELLOW}Please install a LaTeX distribution (e.g., MiKTeX on Windows, TeX Live on Linux).{RESET}")
        return False
    except Exception as e:
        print(f"{RED}Unexpected error during compilation: {e}{RESET}")
        return False
    
def get_and_create_next_run_dir(base_dir: Path) -> Path:
    """
    Finds the next available indexed directory (e.g., 'Results/1')
    and creates it.
    Returns the Path to the newly created directory.
    """
    index = 1
    if not base_dir.is_dir():
        print(f"{RED}Error: The base directory '{base_dir}' does not exist.{RESET}")
        sys.exit(1) 

    while True:
        new_dir_path = base_dir / str(index) 
        
        if not new_dir_path.exists():
            break 
        
        index += 1

    try:
        new_dir_path.mkdir()
        print(f"{GREEN}Successfully created new run directory: {new_dir_path}{RESET}")
        return new_dir_path 
    except OSError as e:
        print(f"{RED}ERROR: Could not create directory: {new_dir_path}{RESET}")
        print(f"{RED}Details: {e}{RESET}")
        sys.exit(1)



def resolve_absolute_bbox(child_bbox, parent):
    """Berechnet die echte Bounding Box anhand der Gruppen-Metriken."""
    scale_x = parent['width'] / parent['chExt_x'] if parent['chExt_x'] else 1
    scale_y = parent['height'] / parent['chExt_y'] if parent['chExt_y'] else 1
    
    # Angenommen, deine bbox nutzt 'l' (left) und 't' (top). Ggf. an Docling-Keys anpassen.
    abs_l = parent['left'] + (child_bbox['l'] - parent['chOff_x']) * scale_x
    abs_t = parent['top'] + (child_bbox['t'] - parent['chOff_y']) * scale_y
    
    child_w = child_bbox['r'] - child_bbox['l']
    child_h = child_bbox['b'] - child_bbox['t']
    
    abs_w = child_w * scale_x
    abs_h = child_h * scale_y
    
    return {
        'l': abs_l, 
        't': abs_t, 
        'r': abs_l + abs_w, 
        'b': abs_t + abs_h
    }

def calculate_geometry(bbox: dict, page_width: float, page_height: float) -> dict | None:
    """
    Calculates relative LaTeX coordinates (0.0-1.0).
    Computes normalized x, y, width, and height based on the bounding box and page dimensions,
    ensuring values remain within valid page boundaries.
    """
    if not bbox or page_width == 0 or page_height == 0:
        return None
    
    left = bbox.get('l', 0)
    top = bbox.get('t', 0)
    right = bbox.get('r', 0)
    bottom = bbox.get('b', 0)
    
    width_emu = abs(right - left)    
    height_emu = abs(bottom - top)    
    visual_top = min(top, bottom)
    
    rel_x = left / page_width
    rel_y = visual_top / page_height
    rel_w = width_emu / page_width
    rel_h = height_emu / page_height

    return {
        "x": round(max(0.0, min(1.0, rel_x)), 3),
        "y": round(max(0.0, min(1.0, rel_y)), 3),
        "w": round(max(0.0, min(1.0, rel_w)), 3),
        "h": round(max(0.0, min(1.0, rel_h)), 3) 
    }

def is_code_line(line: str) -> bool:
    """
    Determines if a text line likely contains programming code by checking for common syntax tokens.
    """
    code_tokens = [';', '{', '}', 'int ', 'public ', 'private ', '=', 'while ', 'if ', 'for ']
    return any(token in line for token in code_tokens)

def build_geo_dict(elements: list) -> dict:
    """
    Groups elements based on their geometry to identify spatially identical items.
    Returns a dictionary mapping geometry tuples to lists of matching elements.
    """
    geos = {}
    for i, el in enumerate(elements):
        geo = tuple(sorted(el['geometry'].items()))
        geos.setdefault(geo, []).append((i, el))
    return geos

def group_elements(elements: list) -> list:
    """
    Groups elements by spatial geometry and semantic type.

    Identifies headers and footers based on vertical position. Detects and merges 
    consecutive code lines into unified code blocks. Classifies remaining text groups 
    as lists or plain text paragraphs based on item count and length, aggregating 
    their geometries and text content accordingly.
    """
    grouped = []
    used = set()
    geos = build_geo_dict(elements)
    
    for geo, group in geos.items():
        first_el = group[0][1] 
        y = first_el['geometry']['y']
        
        if y < 0.03:
            items = [(idx, el) for idx, el in group if 'text' in el and idx not in used]
            if items:
                text = "\n".join(el['text'] for _, el in items)
                grouped.append({
                    "type": "header",
                    "geometry": items[0][1]['geometry'],
                    "text": text.strip(),
                    "fontsize": "3pt", 
                })
                for idx, _ in items: used.add(idx)
        

        elif y > 0.87:
            items = [(idx, el) for idx, el in group if 'text' in el and idx not in used]
            if items:
                text = "\n".join(el['text'] for _, el in items)
                grouped.append({
                    "type": "footer",
                    "geometry": items[0][1]['geometry'],
                    "text": text.strip(),
                    "fontsize": "3pt",
                })
                for idx, _ in items: used.add(idx)

        sure_code_indices = sorted([
            i for i, (idx, el) in enumerate(group)
            if "text" in el and is_code_line(el['text'])
        ])
        
        if len(sure_code_indices) >= 2:
            blocks = []
            current_block = [sure_code_indices[0]]
            for i in range(1, len(sure_code_indices)):
                if sure_code_indices[i] - sure_code_indices[i-1] > 4: 
                    blocks.append(current_block)
                    current_block = []
                current_block.append(sure_code_indices[i])
            blocks.append(current_block)
            
            for blk in blocks:
                if len(blk) < 2: continue
                
                subset = group[blk[0] : blk[-1] + 1]
                code_text = "\n".join(el['text'] for idx, el in subset if 'text' in el)
                union_geo = subset[0][1]['geometry']
                
                grouped.append({
                    "type": "codeblock",
                    "geometry": union_geo,
                    "text": f"\\begin{{lstlisting}}[language=Java]\n{code_text}\n\\end{{lstlisting}}"
                })
                for idx, el in subset: used.add(idx)


        list_like = [(idx, el) for idx, el in group if idx not in used and (
            el['type'] == "list" or el.get("label") in ("list_item", "paragraph", "text"))]
            
        if list_like:
                group_align = "t"
                for _, el in list_like:
                    if el.get("align") == "b":
                        group_align = "b"
                        break

                items = [el['text'] for idx, el in list_like if 'text' in el]
                union_geo = list_like[0][1]['geometry']
                
               
                is_list = False
                
                if len(items) >= 3:
                    if len(items) > 4:
                        is_list = True
                    elif all(len(i) > 20 for i in items):
                        is_list = True

                if is_list:
                    grouped.append({
                        "type": "list",
                        "geometry": union_geo,
                        "items": items,
                        "align": group_align, 
                        "fontsize": "scriptsize"
                    })
                
                else:
                    full_text = "\n".join(items)
                    
                    font_sz = "tiny" if len(full_text) < 20 else "small"
                        
                    grouped.append({
                        "type": "text", 
                        "geometry": union_geo,
                        "text": full_text, 
                        "align": group_align, 
                        "fontsize": font_sz
                    })
                
                for idx, el in list_like: used.add(idx)

        for idx, el in group:
            if idx not in used:
                grouped.append(el)
                used.add(idx)
                
    return grouped


def detect_header_candidate(slides: list) -> tuple | None:
    """
    Identifies a recurring header element across multiple slides.

    Analyzes text elements to find content that repeats in the same position on at least
    70% of the slides. Returns the text and geometry of the longest recurring candidate
    found, or None if no consistent header is detected.
    """
    counter = Counter()
    texts = {}
    for slide in slides:
        for el in slide['elements']:
            if el['type'] == 'text' or el.get('label') in ['paragraph', 'header', 'footer']:
                key = (el.get('text').strip(), tuple(sorted(el['geometry'].items())))
                counter[key] += 1
                texts[key] = el.get('text')
    thresh = int(0.7 * len(slides)) 
    if not counter:
        return None
    candidates = [key for key, val in counter.items() if val >= thresh]
    if not candidates:
        return None
    header_key = max(candidates, key=lambda k: len(k[0]))
    header_text = header_key[0]
    header_geometry = header_key[1]
    return header_text, dict(header_geometry)



def load_slides(json_path: str | Path) -> list | None:
    """
    Loads slide data from a JSON file.

    Verifies the file's existence and attempts to parse it using UTF-8 encoding.
    Returns the parsed JSON data (typically a list of slides) or None if the file
    is missing or corrupt, logging errors with colored output.
    """
    path_obj = Path(json_path)
    
    if not path_obj.exists():
        print(f"{RED}[ERROR] JSON not found at {path_obj}{RESET}")
        return None
    
    try:
        with open(path_obj, 'r', encoding='utf-8') as f:
            return json.load(f)
            
    except json.JSONDecodeError as e:
        print(f"{RED}[ERROR] JSON file is corrupted: {e}{RESET}")
        return None
    
def get_slide_dimensions(pptx_path: str) -> tuple[int, int]:
    """
    Retrieves the total width and height of the presentation slides in EMUs (English Metric Units).
    Returns (0, 0) and logs a warning if the file cannot be read.
    """
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        return prs.slide_width, prs.slide_height
    except Exception as e:
        print(f"{YELLOW}Could not load PPTX dimensions: {e}{RESET}")
        return 0, 0
    

def inject_group_metrics(docling_slides: list, pptx_path: str) -> list:
    """
    Gleicht Docling-Elemente mit python-pptx ab und injiziert 'parent_group',
    falls ein Element Teil einer Gruppe ist.
    """
    prs = Presentation(pptx_path)
    
    for slide_idx, slide_data in enumerate(docling_slides):
        if slide_idx >= len(prs.slides):
            break
            
        pptx_slide = prs.slides[slide_idx]
        docling_elements = slide_data.get('elements', [])
        
        for shape in pptx_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                xfrm = shape._element.grpSpPr.xfrm
                parent_metrics = {
                    'left': shape.left, 'top': shape.top, 
                    'width': shape.width, 'height': shape.height,
                    'chOff_x': xfrm.chOff.x, 'chOff_y': xfrm.chOff.y,
                    'chExt_x': xfrm.chExt.cx, 'chExt_y': xfrm.chExt.cy
                }
                
                for child in shape.shapes:
                    # Finde das passende Docling-Element anhand der fehlerhaften lokalen X-Koordinate
                    # (Toleranzwert einbauen, da Docling intern leicht runden könnte)
                    for el in docling_elements:
                        if 'bbox' in el and abs(el['bbox']['l'] - child.left) < 100:
                            el['parent_group'] = parent_metrics
                            
    return docling_slides

def enrich_and_group_slides(slides: list, slide_width: int, slide_height: int) -> list:
    """
    Processes raw slide data by normalizing element geometry and grouping related items.

    Iterates through every element to convert absolute bounding boxes into relative 
    LaTeX coordinates (0.0-1.0). Afterwards, it reorganizes the elements into semantic 
    groups (headers, code blocks, lists) using the grouping logic.
    """
    for slide in slides:
        elements = slide.get('elements', [])
        
        for el in elements:
            if 'bbox' in el:
                if 'parent_group' in el: 
                    el['bbox'] = resolve_absolute_bbox(el['bbox'], el['parent_group'])
                
                geo = calculate_geometry(el['bbox'], slide_width, slide_height)
                el['geometry'] = geo
                del el['bbox']
            
            if 'parent_group' in el:
                del el['parent_group']
        
        elements = [el for el in elements if not (el.get('type') == 'picture' and 'label' not in el)]
        
        slide['elements'] = group_elements(elements)
        
    return slides

def save_json(data: dict | list, path: str | Path) -> None:
    """
    Saves data to a JSON file with pretty printing.
    
    Writes the provided dictionary or list to the specified path using UTF-8 encoding,
    formatting it with an indentation of 2 spaces for readability and ensuring 
    non-ASCII characters are preserved.
    """
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def sanitize_latex(llm_text: str) -> str:
    """
    Cleans up malformed LaTeX commands often generated by LLMs.

    Normalizes backslashes by replacing control characters, forward slashes, 
    or double backslashes before common keywords (begin, end, item) with a 
    single correct LaTeX backslash.
    """
    latex = re.sub(r'([\x00-\x1F]|\/)+begin', r'\\begin', llm_text, flags=re.MULTILINE)
    latex = re.sub(r'([\x00-\x1F]|\/)+end', r'\\end', latex, flags=re.MULTILINE)

    latex = re.sub(r'([\x00-\x1F]|\/)+item', r'\\item', latex, flags=re.MULTILINE)

    latex = latex.replace('\x08', '\\')

    latex = re.sub(r'\\\\begin', r'\\begin', latex)
    latex = re.sub(r'\\\\end', r'\\end', latex)
    latex = re.sub(r'\\\\item', r'\\item', latex)

    return latex


def repair_latex_output(latex_code: str) -> str:
    """
    Repairs common hallucinated syntax errors in LLM-generated LaTeX.
    
    Specifically targets known issues where models output truncated variable names,
    such as correcting '\paper' to '\paperheight' in geometry definitions.
    """
    latex_code = re.sub(r'\\paper(?!(height|width))', r'\\paperheight', latex_code)
    
    latex_code = latex_code.replace(r'\paper]', r'\paperheight]')
    
    
    return latex_code

from pptx import Presentation

def get_text_alignment_map(pptx_path: str) -> dict:
    """
    Scans the presentation for text boxes that simulate bottom alignment via empty lines.

    Identifies text shapes where the user has pressed 'Enter' multiple times (>= 2 empty paragraphs)
    at the start to push text down. Returns a map of these texts to force 'bottom' alignment
    in the subsequent conversion process.
    """
    prs = Presentation(pptx_path)
    override_map = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_map = {}
        
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text.strip():
                continue
            
            tf = shape.text_frame
            leading_empty_paragraphs = 0
            
            for p in tf.paragraphs:
                if not p.text.strip():
                    leading_empty_paragraphs += 1
                else:
                    break
            
            if leading_empty_paragraphs >= 2:
                clean_key = "".join(shape.text.split()).lower()[:50]
                
                if clean_key:
                    slide_map[clean_key] = "b"

        if slide_map:
            override_map[slide_idx + 1] = slide_map
            
    return override_map