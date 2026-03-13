from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_cm(emu):
    return round(emu / 360000, 2)

def debug_pptx_groups(path):
    prs = Presentation(path)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                g_left, g_top = shape.left, shape.top
                g_width, g_height = shape.width, shape.height
                
                xfrm = shape._element.grpSpPr.xfrm
                ch_off_x, ch_off_y = xfrm.chOff.x, xfrm.chOff.y
                ch_ext_x, ch_ext_y = xfrm.chExt.cx, xfrm.chExt.cy
                
                print(f"GRUPPE: '{shape.name}'")
                print(f"  Folie-Pos (Left/Top): {emu_to_cm(g_left)}cm / {emu_to_cm(g_top)}cm")
                print(f"  Interner Offset (chOff X/Y): {emu_to_cm(ch_off_x)}cm / {emu_to_cm(ch_off_y)}cm")
                print(f"  Interne Größe (chExt W/H): {emu_to_cm(ch_ext_x)}cm / {emu_to_cm(ch_ext_y)}cm")
                
                print("  KINDER dieser Gruppe:")
                for child in shape.shapes:
                    c_left, c_top = child.left, child.top
                    c_width, c_height = child.width, child.height
                    
                    scale_x = g_width / ch_ext_x if ch_ext_x != 0 else 1
                    scale_y = g_height / ch_ext_y if ch_ext_y != 0 else 1
                    
                    abs_left_emu = g_left + (c_left - ch_off_x) * scale_x
                    abs_top_emu = g_top + (c_top - ch_off_y) * scale_y
                    abs_width_emu = c_width * scale_x
                    abs_height_emu = c_height * scale_y
                    
                    print(f"    - '{child.name}':")
                    print(f"      Gelesen (falsch) -> L: {emu_to_cm(c_left)}cm, T: {emu_to_cm(c_top)}cm, W: {emu_to_cm(c_width)}cm, H: {emu_to_cm(c_height)}cm")
                    print(f"      Berechnet (echt) -> L: {emu_to_cm(abs_left_emu)}cm, T: {emu_to_cm(abs_top_emu)}cm, W: {emu_to_cm(abs_width_emu)}cm, H: {emu_to_cm(abs_height_emu)}cm")

# Pfad zu deiner Datei anpassen
debug_pptx_groups("./input/Algorithmik.pptx")