from pptx import Presentation

# Wir behalten den alten Namen 'get_text_alignment_map',
# damit pipeline.py und main.py nicht geändert werden müssen.
def get_text_alignment_map(pptx_path):
    """
    Sucht NUR nach Texten, die mit >= 2 leeren Absätzen (Enters) beginnen.
    Ignoriert PPT-Einstellungen und gibt nur diese 'Fake-Bottom'-Fälle zurück.
    """
    prs = Presentation(pptx_path)
    override_map = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_map = {}
        
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text.strip():
                continue
            
            tf = shape.text_frame
            
            leading_empty_paragraphs = 0
            
            for p in tf.paragraphs:
                if not p.text.strip():
                    leading_empty_paragraphs += 1
                else:
                    break
            
            if leading_empty_paragraphs >= 2:
                clean_key = "".join(shape.text.split()).lower()[:50]
                
                if clean_key:
                    slide_map[clean_key] = "b"

        if slide_map:
            override_map[slide_idx + 1] = slide_map
            
    return override_map
